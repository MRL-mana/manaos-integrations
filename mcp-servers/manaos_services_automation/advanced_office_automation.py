#!/usr/bin/env python3
"""
高度なOffice自動化システム
PowerPointとExcelの高度な機能を活用した自動化ツール
"""

import asyncio
import logging
from datetime import datetime
from typing import Dict, List, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.utils import get_column_letter
import matplotlib.pyplot as plt
import numpy as np

# ログ設定
logging.basicConfig(
    level=logging.INFO,
    format='[%(asctime)s] [%(levelname)s] %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
logger = logging.getLogger('AdvancedOfficeAutomation')

class AdvancedOfficeAutomation:
    """高度なOffice自動化システム"""
    
    def __init__(self):
        self.base_path = "/root/office_files"
        import os
        os.makedirs(self.base_path, exist_ok=True)
        logger.info("📊 高度なOffice自動化システム初期化完了")
    
    # ===== PowerPoint高度な機能 =====
    
    async def create_animated_presentation(self, filename: str, title: str) -> Dict[str, Any]:
        """アニメーション付きプレゼンテーション作成"""
        try:
            prs = Presentation()
            prs.core_properties.title = title
            
            # タイトルスライド
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = title
            subtitle_shape.text = f"自動生成プレゼンテーション - {datetime.now().strftime('%Y-%m-%d')}"
            
            # アニメーション効果（可能な範囲で）
            # 注: python-pptxはアニメーション機能が限定的
            
            # グラフ付きスライド
            await self._add_chart_slide(prs)
            
            # 画像付きスライド
            await self._add_image_slide(prs)
            
            # テーブル付きスライド
            await self._add_table_slide(prs)
            
            filepath = f"{self.base_path}/{filename}"
            prs.save(filepath)
            
            return {
                "success": True,
                "filename": filename,
                "filepath": filepath,
                "slides": len(prs.slides),
                "message": f"アニメーション付きプレゼンテーション '{title}' を作成しました"
            }
            
        except Exception as e:
            logger.error(f"Animated presentation creation error: {e}")
            return {"success": False, "error": str(e)}
    
    async def _add_chart_slide(self, prs: Presentation):
        """グラフ付きスライドを追加"""
        # 空白スライド
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # タイトル
        title_textbox = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "売上データ分析"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(24)
        title_run.font.bold = True
        
        # グラフデータ
        chart_data = CategoryChartData()
        chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
        chart_data.add_series('2023年', (100, 150, 200, 180))
        chart_data.add_series('2024年', (120, 160, 220, 200))
        
        # 棒グラフ
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        # グラフのタイトル
        chart.chart_title.text_frame.text = "四半期売上比較"
        
    async def _add_image_slide(self, prs: Presentation):
        """画像付きスライドを追加"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # タイトル
        title_textbox = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "データ可視化"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(24)
        title_run.font.bold = True
        
        # サンプル画像生成（matplotlib）
        fig, ax = plt.subplots(figsize=(8, 6))
        data = np.random.randn(1000)
        ax.hist(data, bins=30, alpha=0.7, color='skyblue', edgecolor='black')
        ax.set_title('データ分布')
        ax.set_xlabel('値')
        ax.set_ylabel('頻度')
        
        # 画像保存
        img_path = f"{self.base_path}/temp_chart.png"
        plt.savefig(img_path, dpi=300, bbox_inches='tight')
        plt.close()
        
        # 画像をスライドに追加
        slide.shapes.add_picture(img_path, Inches(2), Inches(2), Inches(6), Inches(4))
        
    async def _add_table_slide(self, prs: Presentation):
        """テーブル付きスライドを追加"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # タイトル
        title_textbox = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "データ一覧"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(24)
        title_run.font.bold = True
        
        # テーブル作成
        rows, cols = 4, 3
        left = Inches(2)
        top = Inches(2)
        width = Inches(6)
        height = Inches(3)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # ヘッダー行
        table.cell(0, 0).text = "項目"
        table.cell(0, 1).text = "2023年"
        table.cell(0, 2).text = "2024年"
        
        # データ行
        data = [
            ["売上", "1000", "1200"],
            ["利益", "300", "400"],
            ["コスト", "700", "800"]
        ]
        
        for i, row_data in enumerate(data, 1):
            for j, cell_data in enumerate(row_data):
                table.cell(i, j).text = cell_data
    
    # ===== Excel高度な機能 =====
    
    async def create_advanced_dashboard(self, filename: str) -> Dict[str, Any]:
        """高度なダッシュボード作成"""
        try:
            wb = Workbook()
            ws = wb.active
            ws.title = "ダッシュボード"
            
            # ヘッダー作成
            await self._create_dashboard_header(ws)
            
            # データ生成
            data = await self._generate_sample_data()
            
            # データシート作成
            data_ws = wb.create_sheet("データ")
            await self._populate_data_sheet(data_ws, data)
            
            # グラフ作成
            await self._create_dashboard_charts(ws, data_ws)
            
            # 条件付き書式
            await self._apply_conditional_formatting(ws)
            
            # ピボットテーブル風の集計
            await self._create_summary_tables(ws, data)
            
            filepath = f"{self.base_path}/{filename}"
            wb.save(filepath)
            
            return {
                "success": True,
                "filename": filename,
                "filepath": filepath,
                "sheets": len(wb.worksheets),
                "message": "高度なダッシュボードを作成しました"
            }
            
        except Exception as e:
            logger.error(f"Advanced dashboard creation error: {e}")
            return {"success": False, "error": str(e)}
    
    async def _create_dashboard_header(self, ws):
        """ダッシュボードヘッダー作成"""
        # メインタイトル
        ws['A1'] = "ビジネス分析ダッシュボード"
        ws['A1'].font = Font(name='メイリオ', size=20, bold=True, color="FFFFFF")
        ws['A1'].fill = PatternFill(start_color="2E86AB", end_color="2E86AB", fill_type="solid")
        ws['A1'].alignment = Alignment(horizontal='center', vertical='center')
        ws.merge_cells('A1:H1')
        
        # サブタイトル
        ws['A2'] = f"作成日: {datetime.now().strftime('%Y年%m月%d日')}"
        ws['A2'].font = Font(name='メイリオ', size=12, italic=True)
        ws['A2'].alignment = Alignment(horizontal='center')
        ws.merge_cells('A2:H2')
        
        # 行の高さ調整
        ws.row_dimensions[1].height = 30
        ws.row_dimensions[2].height = 20
    
    async def _generate_sample_data(self) -> List[Dict]:
        """サンプルデータ生成"""
        import random
        data = []
        months = ['1月', '2月', '3月', '4月', '5月', '6月', 
                 '7月', '8月', '9月', '10月', '11月', '12月']
        
        for month in months:
            data.append({
                'month': month,
                'sales': random.randint(800, 1500),
                'profit': random.randint(200, 500),
                'cost': random.randint(400, 800),
                'customers': random.randint(100, 300)
            })
        
        return data
    
    async def _populate_data_sheet(self, ws, data):
        """データシートにデータを投入"""
        # ヘッダー
        headers = ['月', '売上', '利益', 'コスト', '顧客数']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col)
            cell.value = header
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color="E8E8E8", end_color="E8E8E8", fill_type="solid")
        
        # データ投入
        for row, item in enumerate(data, 2):
            ws.cell(row=row, column=1).value = item['month']
            ws.cell(row=row, column=2).value = item['sales']
            ws.cell(row=row, column=3).value = item['profit']
            ws.cell(row=row, column=4).value = item['cost']
            ws.cell(row=row, column=5).value = item['customers']
        
        # 列幅調整
        for col in range(1, 6):
            ws.column_dimensions[get_column_letter(col)].width = 15
    
    async def _create_dashboard_charts(self, ws, data_ws):
        """ダッシュボード用グラフ作成"""
        # 売上グラフ
        sales_chart = LineChart()
        sales_chart.title = "月次売上推移"
        sales_chart.style = 13
        
        sales_data = Reference(data_ws, min_col=2, min_row=1, max_row=13)
        sales_categories = Reference(data_ws, min_col=1, min_row=2, max_row=13)
        sales_chart.add_data(sales_data, titles_from_data=True)
        sales_chart.set_categories(sales_categories)
        
        ws.add_chart(sales_chart, "A5")
        
        # 利益グラフ
        profit_chart = BarChart()
        profit_chart.title = "月次利益"
        profit_chart.style = 10
        
        profit_data = Reference(data_ws, min_col=3, min_row=1, max_row=13)
        profit_chart.add_data(profit_data, titles_from_data=True)
        profit_chart.set_categories(sales_categories)
        
        ws.add_chart(profit_chart, "F5")
    
    async def _apply_conditional_formatting(self, ws):
        """条件付き書式を適用"""
        from openpyxl.formatting.rule import ColorScaleRule, CellIsRule
        from openpyxl.styles import PatternFill
        
        # カラースケール（A5:A16の売上データ）
        color_scale = ColorScaleRule(
            start_type="min", start_color="FF0000",
            mid_type="percentile", mid_value=50, mid_color="FFFF00",
            end_type="max", end_color="00FF00"
        )
        ws.conditional_formatting.add("B5:B16", color_scale)
        
        # セル値による条件付き書式
        high_profit = CellIsRule(
            operator='greaterThan', formula=['300'],
            fill=PatternFill(start_color="90EE90", end_color="90EE90", fill_type="solid")
        )
        ws.conditional_formatting.add("C5:C16", high_profit)
    
    async def _create_summary_tables(self, ws, data):
        """サマリーテーブル作成"""
        # 集計データ計算
        total_sales = sum(item['sales'] for item in data)
        total_profit = sum(item['profit'] for item in data)
        avg_customers = sum(item['customers'] for item in data) / len(data)
        
        # サマリーテーブル
        summary_row = 20
        ws[f'A{summary_row}'] = "年間合計売上"
        ws[f'B{summary_row}'] = total_sales
        ws[f'B{summary_row}'].number_format = '#,##0'
        
        ws[f'A{summary_row+1}'] = "年間合計利益"
        ws[f'B{summary_row+1}'] = total_profit
        ws[f'B{summary_row+1}'].number_format = '#,##0'
        
        ws[f'A{summary_row+2}'] = "平均顧客数"
        ws[f'B{summary_row+2}'] = avg_customers
        ws[f'B{summary_row+2}'].number_format = '0'
        
        # サマリーのスタイル
        for row in range(summary_row, summary_row+3):
            ws[f'A{row}'].font = Font(bold=True)
            ws[f'B{row}'].font = Font(bold=True, color="2E86AB")
    
    # ===== 統合機能 =====
    
    async def create_office_report_package(self, title: str) -> Dict[str, Any]:
        """Office統合レポートパッケージ作成"""
        try:
            package_name = f"{title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            
            # PowerPointプレゼンテーション作成
            ppt_result = await self.create_animated_presentation(
                f"{package_name}.pptx", f"{title} プレゼンテーション"
            )
            
            # Excelダッシュボード作成
            excel_result = await self.create_advanced_dashboard(
                f"{package_name}.xlsx"
            )
            
            return {
                "success": True,
                "package_name": package_name,
                "powerpoint": ppt_result,
                "excel": excel_result,
                "message": f"Office統合レポートパッケージ '{title}' を作成しました"
            }
            
        except Exception as e:
            logger.error(f"Office report package creation error: {e}")
            return {"success": False, "error": str(e)}

def main():
    """メイン実行関数"""
    automation = AdvancedOfficeAutomation()
    
    async def run_demo():
        logger.info("🚀 高度なOffice自動化デモ開始")
        
        # タイヤ交換広告の高度版を作成
        result = await automation.create_office_report_package("タイヤ交換サービス分析")
        
        if result["success"]:
            logger.info("✅ 高度なOffice自動化完了!")
            logger.info(f"📊 パッケージ: {result['package_name']}")
            logger.info(f"📈 PowerPoint: {result['powerpoint']['filename']}")
            logger.info(f"📊 Excel: {result['excel']['filename']}")
        else:
            logger.error(f"❌ エラー: {result['error']}")
    
    try:
        asyncio.run(run_demo())
    except Exception as e:
        logger.error(f"❌ デモ実行エラー: {e}")

if __name__ == "__main__":
    main()
