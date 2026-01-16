#!/usr/bin/env python3
"""
高度なOffice MCP Server
PowerPointとExcelの高度な機能を統合したMCPサーバー
"""

import asyncio
import logging
import os
from datetime import datetime
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import LineChart, Reference
from openpyxl.utils import get_column_letter
from openpyxl.formatting.rule import ColorScaleRule, CellIsRule

# ログ設定
logging.basicConfig(
    level=logging.INFO,
    format='[%(asctime)s] [%(levelname)s] %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
logger = logging.getLogger('AdvancedOfficeMCP')

class AdvancedOfficeMCPServer:
    """高度なOffice MCP Server"""
    
    def __init__(self, port=int(os.getenv("PORT", "5030"))):
        self.port = port
        self.base_path = os.path.join(os.getenv("HOME", "/root"), "office_files")
        os.makedirs(self.base_path, exist_ok=True)
        logger.info(f"📊 高度なOffice MCP Server 初期化完了 (ポート: {port})")
    
    # ===== PowerPoint高度機能 =====
    
    async def create_business_presentation(self, filename: str, title: str, data: Dict[str, Any]) -> Dict[str, Any]:
        """ビジネスプレゼンテーション作成"""
        try:
            prs = Presentation()
            prs.core_properties.title = title
            
            # タイトルスライド
            await self._create_title_slide(prs, title)
            
            # 概要スライド
            await self._create_overview_slide(prs, data)
            
            # データ分析スライド
            await self._create_data_analysis_slides(prs, data)
            
            # 結論スライド
            await self._create_conclusion_slide(prs, data)
            
            filepath = f"{self.base_path}/{filename}"
            prs.save(filepath)
            
            return {
                "success": True,
                "filename": filename,
                "filepath": filepath,
                "slides": len(prs.slides),
                "message": f"ビジネスプレゼンテーション '{title}' を作成しました"
            }
            
        except Exception as e:
            logger.error(f"Business presentation creation error: {e}")
            return {"success": False, "error": str(e)}
    
    async def _create_title_slide(self, prs: Presentation, title: str):
        """タイトルスライド作成"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = f"自動生成プレゼンテーション\n{datetime.now().strftime('%Y年%m月%d日')}"
        
        # スタイル調整
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 138)
        
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    
    async def _create_overview_slide(self, prs: Presentation, data: Dict[str, Any]):
        """概要スライド作成"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # タイトル
        title_textbox = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "概要"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(32)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(30, 58, 138)
        
        # 概要内容
        overview_textbox = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
        overview_frame = overview_textbox.text_frame
        overview_frame.text = "• データ分析結果の報告\n• 主要指標の動向\n• 改善提案とアクションプラン\n• 今後の展望"
        
        for paragraph in overview_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = RGBColor(50, 50, 50)
    
    async def _create_data_analysis_slides(self, prs: Presentation, data: Dict[str, Any]):
        """データ分析スライド作成"""
        # グラフスライド
        await self._create_chart_slide(prs, data)
        
        # テーブルスライド
        await self._create_table_slide(prs, data)
    
    async def _create_chart_slide(self, prs: Presentation, data: Dict[str, Any]):
        """グラフスライド作成"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # タイトル
        title_textbox = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "データ分析 - 売上推移"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(28)
        title_run.font.bold = True
        
        # サンプルデータでグラフ作成
        chart_data = CategoryChartData()
        chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
        chart_data.add_series('2023年', (100, 150, 200, 180))
        chart_data.add_series('2024年', (120, 160, 220, 200))
        
        # 棒グラフ
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        chart.chart_title.text_frame.text = "四半期売上比較"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
    
    async def _create_table_slide(self, prs: Presentation, data: Dict[str, Any]):
        """テーブルスライド作成"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # タイトル
        title_textbox = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "詳細データ"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(28)
        title_run.font.bold = True
        
        # テーブル作成
        rows, cols = 5, 4
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        height = Inches(3)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # ヘッダー
        headers = ["項目", "2023年", "2024年", "増減率"]
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 58, 138)
        
        # データ行
        table_data = [
            ["売上", "1,000", "1,200", "20%"],
            ["利益", "300", "400", "33%"],
            ["コスト", "700", "800", "14%"],
            ["顧客数", "500", "600", "20%"]
        ]
        
        for row, row_data in enumerate(table_data, 1):
            for col, cell_data in enumerate(row_data):
                cell = table.cell(row, col)
                cell.text = cell_data
                cell.text_frame.paragraphs[0].font.size = Pt(14)
    
    async def _create_conclusion_slide(self, prs: Presentation, data: Dict[str, Any]):
        """結論スライド作成"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # タイトル
        title_textbox = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "結論とアクションプラン"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(32)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(30, 58, 138)
        
        # 結論内容
        conclusion_textbox = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
        conclusion_frame = conclusion_textbox.text_frame
        conclusion_frame.text = "• 売上は前年比20%増加\n• 利益率の改善が顕著\n• 顧客数の拡大が継続\n• 今後も成長戦略を継続"
        
        for paragraph in conclusion_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = RGBColor(50, 50, 50)
    
    # ===== Excel高度機能 =====
    
    async def create_financial_dashboard(self, filename: str, data: Dict[str, Any]) -> Dict[str, Any]:
        """財務ダッシュボード作成"""
        try:
            wb = Workbook()
            ws = wb.active
            ws.title = "財務ダッシュボード"
            
            # ダッシュボードヘッダー
            await self._create_financial_header(ws)
            
            # KPI表示エリア
            await self._create_kpi_area(ws, data)
            
            # グラフエリア
            await self._create_financial_charts(ws, data)
            
            # 詳細データシート
            detail_ws = wb.create_sheet("詳細データ")
            await self._populate_financial_data(detail_ws, data)
            
            # 条件付き書式
            await self._apply_financial_formatting(ws)
            
            filepath = f"{self.base_path}/{filename}"
            wb.save(filepath)
            
            return {
                "success": True,
                "filename": filename,
                "filepath": filepath,
                "sheets": len(wb.worksheets),
                "message": "財務ダッシュボードを作成しました"
            }
            
        except Exception as e:
            logger.error(f"Financial dashboard creation error: {e}")
            return {"success": False, "error": str(e)}
    
    async def _create_financial_header(self, ws):
        """財務ダッシュボードヘッダー"""
        # メインタイトル
        ws['A1'] = "財務ダッシュボード"
        ws['A1'].font = Font(name='メイリオ', size=24, bold=True, color="FFFFFF")
        ws['A1'].fill = PatternFill(start_color="1E3A8A", end_color="1E3A8A", fill_type="solid")
        ws['A1'].alignment = Alignment(horizontal='center', vertical='center')
        ws.merge_cells('A1:H1')
        
        # サブタイトル
        ws['A2'] = f"更新日: {datetime.now().strftime('%Y年%m月%d日 %H:%M')}"
        ws['A2'].font = Font(name='メイリオ', size=12, italic=True)
        ws['A2'].alignment = Alignment(horizontal='center')
        ws.merge_cells('A2:H2')
        
        # 行の高さ調整
        ws.row_dimensions[1].height = 40
        ws.row_dimensions[2].height = 25
    
    async def _create_kpi_area(self, ws, data: Dict[str, Any]):
        """KPI表示エリア"""
        # KPIタイトル
        ws['A4'] = "主要KPI"
        ws['A4'].font = Font(name='メイリオ', size=16, bold=True)
        ws.merge_cells('A4:H4')
        
        # KPIカード
        kpis = [
            ("売上", "1,200,000", "円", "20%"),
            ("利益", "400,000", "円", "33%"),
            ("利益率", "33.3", "%", "2.1%"),
            ("顧客数", "600", "名", "20%")
        ]
        
        for i, (name, value, unit, change) in enumerate(kpis):
            row = 6 + (i // 2) * 3
            col = 1 + (i % 2) * 4
            
            # KPIカード背景
            ws.cell(row=row, column=col, value=name).font = Font(bold=True, size=12)
            ws.cell(row=row+1, column=col, value=f"{value} {unit}").font = Font(bold=True, size=16, color="1E3A8A")
            ws.cell(row=row+2, column=col, value=f"前年比: {change}").font = Font(size=10, color="228B22")
            
            # カード背景色
            for r in range(row, row+3):
                for c in range(col, col+3):
                    cell = ws.cell(row=r, column=c)
                    cell.fill = PatternFill(start_color="F8F9FA", end_color="F8F9FA", fill_type="solid")
                    cell.border = Border(
                        left=Side(style='thin'),
                        right=Side(style='thin'),
                        top=Side(style='thin'),
                        bottom=Side(style='thin')
                    )
    
    async def _create_financial_charts(self, ws, data: Dict[str, Any]):
        """財務グラフ作成"""
        # 売上グラフ
        sales_chart = LineChart()
        sales_chart.title = "月次売上推移"
        sales_chart.style = 13
        sales_chart.y_axis.title = "売上 (千円)"
        sales_chart.x_axis.title = "月"
        
        # サンプルデータ
        months = ['1月', '2月', '3月', '4月', '5月', '6月', '7月', '8月', '9月', '10月', '11月', '12月']
        sales_data = [80, 90, 100, 95, 110, 120, 115, 125, 130, 135, 140, 145]
        
        # データシートにデータを追加
        data_ws = ws.parent.create_sheet("グラフデータ")
        data_ws['A1'] = "月"
        data_ws['B1'] = "売上"
        
        for i, (month, sales) in enumerate(zip(months, sales_data), 2):
            data_ws[f'A{i}'] = month
            data_ws[f'B{i}'] = sales
        
        # グラフデータ参照
        sales_ref = Reference(data_ws, min_col=2, min_row=1, max_row=13)
        categories_ref = Reference(data_ws, min_col=1, min_row=2, max_row=13)
        sales_chart.add_data(sales_ref, titles_from_data=True)
        sales_chart.set_categories(categories_ref)
        
        ws.add_chart(sales_chart, "A15")
    
    async def _populate_financial_data(self, ws, data: Dict[str, Any]):
        """財務データ投入"""
        # ヘッダー
        headers = ["月", "売上", "利益", "コスト", "顧客数", "利益率"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col)
            cell.value = header
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color="1E3A8A", end_color="1E3A8A", fill_type="solid")
            cell.font = Font(bold=True, color="FFFFFF")
        
        # サンプルデータ
        months = ['1月', '2月', '3月', '4月', '5月', '6月', '7月', '8月', '9月', '10月', '11月', '12月']
        sales_data = [800, 900, 1000, 950, 1100, 1200, 1150, 1250, 1300, 1350, 1400, 1450]
        profit_data = [240, 270, 300, 285, 330, 360, 345, 375, 390, 405, 420, 435]
        cost_data = [560, 630, 700, 665, 770, 840, 805, 875, 910, 945, 980, 1015]
        customer_data = [100, 110, 120, 115, 130, 140, 135, 145, 150, 155, 160, 165]
        
        for i, (month, sales, profit, cost, customers) in enumerate(zip(months, sales_data, profit_data, cost_data, customer_data), 2):
            ws.cell(row=i, column=1).value = month
            ws.cell(row=i, column=2).value = sales
            ws.cell(row=i, column=3).value = profit
            ws.cell(row=i, column=4).value = cost
            ws.cell(row=i, column=5).value = customers
            ws.cell(row=i, column=6).value = f"{profit/sales*100:.1f}%"
        
        # 列幅調整
        for col in range(1, 7):
            ws.column_dimensions[get_column_letter(col)].width = 15
    
    async def _apply_financial_formatting(self, ws):
        """財務書式設定"""
        # 売上データのカラースケール
        sales_range = "B15:B26"
        color_scale = ColorScaleRule(
            start_type="min", start_color="FF0000",
            mid_type="percentile", mid_value=50, mid_color="FFFF00",
            end_type="max", end_color="00FF00"
        )
        ws.conditional_formatting.add(sales_range, color_scale)
        
        # 利益率の条件付き書式
        profit_rate_rule = CellIsRule(
            operator='greaterThan', formula=['30'],
            fill=PatternFill(start_color="90EE90", end_color="90EE90", fill_type="solid")
        )
        ws.conditional_formatting.add("F15:F26", profit_rate_rule)
    
    # ===== 統合機能 =====
    
    async def create_comprehensive_report(self, title: str, data: Dict[str, Any]) -> Dict[str, Any]:
        """包括的レポート作成"""
        try:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            package_name = f"{title}_{timestamp}"
            
            # PowerPointビジネスプレゼンテーション
            ppt_result = await self.create_business_presentation(
                f"{package_name}.pptx", f"{title} プレゼンテーション", data
            )
            
            # Excel財務ダッシュボード
            excel_result = await self.create_financial_dashboard(
                f"{package_name}.xlsx", data
            )
            
            return {
                "success": True,
                "package_name": package_name,
                "powerpoint": ppt_result,
                "excel": excel_result,
                "message": f"包括的レポートパッケージ '{title}' を作成しました"
            }
            
        except Exception as e:
            logger.error(f"Comprehensive report creation error: {e}")
            return {"success": False, "error": str(e)}

def main():
    """メイン実行関数"""
    server = AdvancedOfficeMCPServer()
    
    async def run_demo():
        logger.info("🚀 高度なOffice MCP Server デモ開始")
        
        # タイヤ交換サービス分析レポート作成
        sample_data = {
            "company": "タイヤ交換サービス",
            "period": "2024年度",
            "metrics": {
                "revenue": 1200000,
                "profit": 400000,
                "customers": 600
            }
        }
        
        result = await server.create_comprehensive_report("タイヤ交換サービス分析", sample_data)
        
        if result["success"]:
            logger.info("✅ 高度なOffice MCP Server デモ完了!")
            logger.info(f"📊 パッケージ: {result['package_name']}")
            logger.info(f"📈 PowerPoint: {result['powerpoint']['slides']}スライド")
            logger.info(f"📊 Excel: {result['excel']['sheets']}シート")
        else:
            logger.error(f"❌ エラー: {result['error']}")
    
    try:
        asyncio.run(run_demo())
    except KeyboardInterrupt:
        logger.info("🛑 高度なOffice MCP Server 停止")
    except Exception as e:
        logger.error(f"❌ 高度なOffice MCP Server エラー: {e}")

if __name__ == "__main__":
    main()
