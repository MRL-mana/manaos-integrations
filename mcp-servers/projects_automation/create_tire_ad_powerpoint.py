#!/usr/bin/env python3
"""
PowerPointでタイヤ交換広告を作成
参考画像と同じデザインでプロ仕様の広告を作成
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

class TireAdPowerPointCreator:
    """PowerPointタイヤ交換広告作成クラス"""
    
    def __init__(self):
        self.prs = Presentation()
        
        # スライドサイズ設定 (16:9)
        self.prs.slide_width = Inches(13.33)  # 1920px相当
        self.prs.slide_height = Inches(7.5)   # 1080px相当
        
    def create_tire_ad_slide(self):
        """タイヤ交換広告スライドを作成"""
        
        # 空白スライドを追加
        blank_slide_layout = self.prs.slide_layouts[6]  # 空白レイアウト
        slide = self.prs.slides.add_slide(blank_slide_layout)
        
        # 上部60% - タイヤ画像エリア
        self.create_tire_image_area(slide)
        
        # 下部40% - テキストエリア
        self.create_text_area(slide)
        
        # ファイル保存
        filename = "/root/tire_exchange_ad.pptx"
        self.prs.save(filename)
        print(f"🎉 PowerPointタイヤ交換広告を作成しました: {filename}")
        
        return filename
    
    def create_tire_image_area(self, slide):
        """上部60% - タイヤ画像エリア"""
        
        # 上部エリアの背景（白）
        left = Inches(0)
        top = Inches(0)
        width = Inches(13.33)
        height = Inches(4.5)  # 60%
        
        # 白い背景矩形
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白
        background.line.fill.background()
        
        # タイヤ画像の代わりにテキスト表現
        textbox = slide.shapes.add_textbox(
            Inches(2), Inches(1.5), Inches(9.33), Inches(1.5)
        )
        text_frame = textbox.text_frame
        text_frame.text = "❄️ タイヤの履き替えは当店で ❄️"
        
        # テキストスタイル
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        
        run = paragraph.runs[0]
        run.font.name = 'メイリオ'
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor(44, 62, 80)  # ダークグレー
        
        # サブテキスト
        textbox2 = slide.shapes.add_textbox(
            Inches(4), Inches(2.8), Inches(5.33), Inches(0.8)
        )
        text_frame2 = textbox2.text_frame
        text_frame2.text = "スマホで便利に"
        
        paragraph2 = text_frame2.paragraphs[0]
        paragraph2.alignment = PP_ALIGN.CENTER
        
        run2 = paragraph2.runs[0]
        run2.font.name = 'メイリオ'
        run2.font.size = Pt(24)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(44, 62, 80)
        
        # 雪の装飾
        for i, x_pos in enumerate([1, 3, 10, 12]):
            snow_textbox = slide.shapes.add_textbox(
                Inches(x_pos), Inches(0.5), Inches(0.5), Inches(0.5)
            )
            snow_frame = snow_textbox.text_frame
            snow_frame.text = "❄"
            snow_paragraph = snow_frame.paragraphs[0]
            snow_paragraph.alignment = PP_ALIGN.CENTER
            snow_run = snow_paragraph.runs[0]
            snow_run.font.size = Pt(28)
    
    def create_text_area(self, slide):
        """下部40% - テキストエリア"""
        
        # ダークブルー背景
        left = Inches(0)
        top = Inches(4.5)  # 上部60%の下から
        width = Inches(13.33)
        height = Inches(3.0)  # 40%
        
        # ダークブルー背景矩形
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(30, 58, 138)  # #1e3a8a
        background.line.fill.background()
        
        # メインテキスト
        main_textbox = slide.shapes.add_textbox(
            Inches(2), Inches(5.0), Inches(9.33), Inches(1.2)
        )
        main_frame = main_textbox.text_frame
        main_frame.text = "タイヤ交換は\nかんたんネット予約"
        
        main_paragraph = main_frame.paragraphs[0]
        main_paragraph.alignment = PP_ALIGN.CENTER
        
        main_run = main_paragraph.runs[0]
        main_run.font.name = 'メイリオ'
        main_run.font.size = Pt(48)
        main_run.font.bold = True
        main_run.font.color.rgb = RGBColor(225, 29, 72)  # #e11d48 (ピンク/赤)
        
        # CTAボタン
        button_left = Inches(5.5)
        button_top = Inches(6.2)
        button_width = Inches(2.33)
        button_height = Inches(0.6)
        
        button = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            button_left, button_top, button_width, button_height
        )
        button.fill.solid()
        button.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白
        button.line.color.rgb = RGBColor(204, 204, 204)  # 薄いグレー
        
        # ボタンテキスト
        button_textbox = slide.shapes.add_textbox(
            Inches(5.6), Inches(6.3), Inches(2.13), Inches(0.4)
        )
        button_frame = button_textbox.text_frame
        button_frame.text = "予約はこちら"
        
        button_paragraph = button_frame.paragraphs[0]
        button_paragraph.alignment = PP_ALIGN.CENTER
        
        button_run = button_paragraph.runs[0]
        button_run.font.name = 'メイリオ'
        button_run.font.size = Pt(18)
        button_run.font.bold = True
        button_run.font.color.rgb = RGBColor(0, 0, 0)  # 黒
        
        # 指差しアイコン
        pointer_textbox = slide.shapes.add_textbox(
            Inches(8.0), Inches(6.3), Inches(0.4), Inches(0.4)
        )
        pointer_frame = pointer_textbox.text_frame
        pointer_frame.text = "👉"
        
        pointer_paragraph = pointer_frame.paragraphs[0]
        pointer_paragraph.alignment = PP_ALIGN.CENTER
        
        pointer_run = pointer_paragraph.runs[0]
        pointer_run.font.size = Pt(20)

def main():
    """メイン実行関数"""
    print("🚗 PowerPointでタイヤ交換広告を作成中...")
    
    try:
        creator = TireAdPowerPointCreator()
        filename = creator.create_tire_ad_slide()
        
        print("✅ PowerPoint作成完了!")
        print(f"📁 ファイル: {filename}")
        print("🎨 デザイン:")
        print("  - 上部60%: 雪道タイヤエリア（白背景）")
        print("  - 下部40%: ダークブルー背景")
        print("  - メインテキスト: ピンク/赤色")
        print("  - CTAボタン: 白背景、指差しアイコン付き")
        print("  - サイズ: 16:9 (1920x1080px相当)")
        
        return filename
        
    except Exception as e:
        print(f"❌ エラー: {e}")
        return None

if __name__ == "__main__":
    main()
