#!/usr/bin/env python3
"""
PowerPoint MCP Server - PowerPoint自動化システム
X280やリモート環境でPowerPointファイルを自動生成・編集
"""

import asyncio
import logging
import os
from datetime import datetime
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ログ設定
logging.basicConfig(
    level=logging.INFO,
    format='[%(asctime)s] [%(levelname)s] %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
logger = logging.getLogger('PowerPointMCP')

class PowerPointMCPServer:
    """PowerPoint MCP Server"""
    
    def __init__(self, port=int(os.getenv("PORT", "5025"))):
        self.port = port
        self.base_path = "/root/powerpoint_files"
        os.makedirs(self.base_path, exist_ok=True)
        logger.info(f"📊 PowerPoint MCP Server 初期化完了 (ポート: {port})")
    
    async def create_presentation(self, filename: str, title: str = "New Presentation") -> Dict[str, Any]:
        """新しいプレゼンテーションを作成"""
        try:
            prs = Presentation()
            prs.core_properties.title = title
            
            # デフォルトスライド追加
            slide_layout = prs.slide_layouts[0]  # タイトルスライド
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = title
            subtitle_shape.text = f"Created by PowerPoint MCP Server - {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            
            # ファイル保存
            filepath = os.path.join(self.base_path, filename)
            prs.save(filepath)
            
            return {
                "success": True,
                "filename": filename,
                "filepath": filepath,
                "message": f"プレゼンテーション '{title}' を作成しました"
            }
            
        except Exception as e:
            logger.error(f"Presentation creation error: {e}")
            return {"success": False, "error": str(e)}
    
    async def add_slide(self, filename: str, slide_type: str = "blank", content: Dict[str, Any] = None) -> Dict[str, Any]:
        """スライドを追加"""
        try:
            filepath = os.path.join(self.base_path, filename)
            if not os.path.exists(filepath):
                return {"success": False, "error": f"ファイルが見つかりません: {filename}"}
            
            prs = Presentation(filepath)
            
            # スライドタイプに応じてレイアウト選択
            layout_map = {
                "blank": 6,
                "title": 0,
                "content": 1,
                "section_header": 2,
                "two_content": 3,
                "comparison": 4,
                "title_only": 5
            }
            
            layout_idx = layout_map.get(slide_type, 6)
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)
            
            # コンテンツ追加
            if content:
                await self._add_content_to_slide(slide, content)
            
            prs.save(filepath)
            
            return {
                "success": True,
                "filename": filename,
                "slide_type": slide_type,
                "message": f"{slide_type}スライドを追加しました"
            }
            
        except Exception as e:
            logger.error(f"Add slide error: {e}")
            return {"success": False, "error": str(e)}
    
    async def create_tire_ad_slide(self, filename: str) -> Dict[str, Any]:
        """タイヤ交換広告スライドを作成"""
        try:
            prs = Presentation()
            
            # スライドサイズ設定 (16:9)
            prs.slide_width = Inches(13.33)  # 1920px相当
            prs.slide_height = Inches(7.5)   # 1080px相当
            
            # 空白スライドを追加
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 上部60% - タイヤ画像エリア
            await self._create_tire_image_area(slide)
            
            # 下部40% - テキストエリア
            await self._create_tire_text_area(slide)
            
            # ファイル保存
            filepath = os.path.join(self.base_path, filename)
            prs.save(filepath)
            
            return {
                "success": True,
                "filename": filename,
                "filepath": filepath,
                "message": "タイヤ交換広告スライドを作成しました"
            }
            
        except Exception as e:
            logger.error(f"Tire ad creation error: {e}")
            return {"success": False, "error": str(e)}
    
    async def _create_tire_image_area(self, slide):
        """上部60% - タイヤ画像エリア"""
        
        # 上部エリアの背景（白）
        left = Inches(0)
        top = Inches(0)
        width = Inches(13.33)
        height = Inches(4.5)  # 60%
        
        # 白い背景矩形
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白
        background.line.fill.background()
        
        # タイヤ画像の代わりにテキスト表現
        textbox = slide.shapes.add_textbox(
            Inches(2), Inches(1.5), Inches(9.33), Inches(1.5)
        )
        text_frame = textbox.text_frame
        text_frame.text = "❄️ タイヤの履き替えは当店で ❄️"
        
        # テキストスタイル
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        
        run = paragraph.runs[0]
        run.font.name = 'メイリオ'
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor(44, 62, 80)  # ダークグレー
        
        # サブテキスト
        textbox2 = slide.shapes.add_textbox(
            Inches(4), Inches(2.8), Inches(5.33), Inches(0.8)
        )
        text_frame2 = textbox2.text_frame
        text_frame2.text = "スマホで便利に"
        
        paragraph2 = text_frame2.paragraphs[0]
        paragraph2.alignment = PP_ALIGN.CENTER
        
        run2 = paragraph2.runs[0]
        run2.font.name = 'メイリオ'
        run2.font.size = Pt(24)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(44, 62, 80)
    
    async def _create_tire_text_area(self, slide):
        """下部40% - テキストエリア"""
        
        # ダークブルー背景
        left = Inches(0)
        top = Inches(4.5)  # 上部60%の下から
        width = Inches(13.33)
        height = Inches(3.0)  # 40%
        
        # ダークブルー背景矩形
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(30, 58, 138)  # #1e3a8a
        background.line.fill.background()
        
        # メインテキスト
        main_textbox = slide.shapes.add_textbox(
            Inches(2), Inches(5.0), Inches(9.33), Inches(1.2)
        )
        main_frame = main_textbox.text_frame
        main_frame.text = "タイヤ交換は\nかんたんネット予約"
        
        main_paragraph = main_frame.paragraphs[0]
        main_paragraph.alignment = PP_ALIGN.CENTER
        
        main_run = main_paragraph.runs[0]
        main_run.font.name = 'メイリオ'
        main_run.font.size = Pt(48)
        main_run.font.bold = True
        main_run.font.color.rgb = RGBColor(225, 29, 72)  # #e11d48 (ピンク/赤)
        
        # CTAボタン
        button_left = Inches(5.5)
        button_top = Inches(6.2)
        button_width = Inches(2.33)
        button_height = Inches(0.6)
        
        button = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            button_left, button_top, button_width, button_height
        )
        button.fill.solid()
        button.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白
        button.line.color.rgb = RGBColor(204, 204, 204)  # 薄いグレー
        
        # ボタンテキスト
        button_textbox = slide.shapes.add_textbox(
            Inches(5.6), Inches(6.3), Inches(2.13), Inches(0.4)
        )
        button_frame = button_textbox.text_frame
        button_frame.text = "予約はこちら"
        
        button_paragraph = button_frame.paragraphs[0]
        button_paragraph.alignment = PP_ALIGN.CENTER
        
        button_run = button_paragraph.runs[0]
        button_run.font.name = 'メイリオ'
        button_run.font.size = Pt(18)
        button_run.font.bold = True
        button_run.font.color.rgb = RGBColor(0, 0, 0)  # 黒
        
        # 指差しアイコン
        pointer_textbox = slide.shapes.add_textbox(
            Inches(8.0), Inches(6.3), Inches(0.4), Inches(0.4)
        )
        pointer_frame = pointer_textbox.text_frame
        pointer_frame.text = "👉"
        
        pointer_paragraph = pointer_frame.paragraphs[0]
        pointer_paragraph.alignment = PP_ALIGN.CENTER
        
        pointer_run = pointer_paragraph.runs[0]
        pointer_run.font.size = Pt(20)
    
    async def _add_content_to_slide(self, slide, content: Dict[str, Any]):
        """スライドにコンテンツを追加"""
        
        if "title" in content:
            # タイトル追加
            title_textbox = slide.shapes.add_textbox(
                Inches(1), Inches(1), Inches(11.33), Inches(1)
            )
            title_frame = title_textbox.text_frame
            title_frame.text = content["title"]
            
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.alignment = PP_ALIGN.CENTER
            
            title_run = title_paragraph.runs[0]
            title_run.font.name = 'メイリオ'
            title_run.font.size = Pt(32)
            title_run.font.bold = True
        
        if "content" in content:
            # コンテンツ追加
            content_textbox = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(11.33), Inches(3)
            )
            content_frame = content_textbox.text_frame
            content_frame.text = content["content"]
            
            content_paragraph = content_frame.paragraphs[0]
            content_paragraph.alignment = PP_ALIGN.LEFT
            
            content_run = content_paragraph.runs[0]
            content_run.font.name = 'メイリオ'
            content_run.font.size = Pt(18)
    
    async def list_presentations(self) -> Dict[str, Any]:
        """プレゼンテーションファイル一覧を取得"""
        try:
            files = []
            for file in os.listdir(self.base_path):
                if file.endswith(('.pptx', '.ppt')):
                    filepath = os.path.join(self.base_path, file)
                    stat = os.stat(filepath)
                    files.append({
                        "filename": file,
                        "size": stat.st_size,
                        "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                        "modified": datetime.fromtimestamp(stat.st_mtime).isoformat()
                    })
            
            return {
                "success": True,
                "files": files,
                "count": len(files)
            }
            
        except Exception as e:
            logger.error(f"List presentations error: {e}")
            return {"success": False, "error": str(e)}
    
    async def get_presentation_info(self, filename: str) -> Dict[str, Any]:
        """プレゼンテーション情報を取得"""
        try:
            filepath = os.path.join(self.base_path, filename)
            if not os.path.exists(filepath):
                return {"success": False, "error": f"ファイルが見つかりません: {filename}"}
            
            prs = Presentation(filepath)
            
            return {
                "success": True,
                "filename": filename,
                "title": prs.core_properties.title,
                "slide_count": len(prs.slides),
                "slide_size": {
                    "width": prs.slide_width,
                    "height": prs.slide_height
                }
            }
            
        except Exception as e:
            logger.error(f"Get presentation info error: {e}")
            return {"success": False, "error": str(e)}

# MCP Server実装
def main():
    """メイン関数"""
    server = PowerPointMCPServer()
    
    async def run_server():
        logger.info("🚀 PowerPoint MCP Server 起動中...")
        
        # テスト用タイヤ交換広告を作成
        result = await server.create_tire_ad_slide("tire_exchange_ad_sample.pptx")
        if result["success"]:
            logger.info(f"📊 サンプル広告作成完了: {result['filename']}")
        
        # サーバーを継続実行
        while True:
            await asyncio.sleep(60)  # 1分ごとにチェック
    
    try:
        asyncio.run(run_server())
    except KeyboardInterrupt:
        logger.info("🛑 PowerPoint MCP Server 停止")
    except Exception as e:
        logger.error(f"❌ PowerPoint MCP Server エラー: {e}")

if __name__ == "__main__":
    main()
